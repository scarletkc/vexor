from pathlib import Path
from types import SimpleNamespace

from docx import Document
from pptx import Presentation

import vexor.services.content_extract_service as ces
from vexor.services.content_extract_service import (
    extract_code_chunks,
    extract_outline_chunks,
    extract_full_chunks,
    extract_full_chunks_with_lines,
    extract_head,
    HEAD_CHAR_LIMIT,
)


def test_extract_head_from_text(tmp_path):
    text_file = tmp_path / "sample.txt"
    text_file.write_text("Title\n\nLine one\nLine two\n")

    snippet = extract_head(text_file, char_limit=HEAD_CHAR_LIMIT)

    assert snippet == "Title Line one Line two"


def test_extract_head_unknown_extension(tmp_path):
    data_file = tmp_path / "data.bin"
    data_file.write_bytes(b"\x00\x01\x02")

    assert extract_head(data_file) is None


def test_extract_head_from_docx(tmp_path):
    doc_path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Docx Title")
    document.add_paragraph("First paragraph")
    document.save(doc_path)

    snippet = extract_head(doc_path)

    assert snippet.startswith("Docx Title First paragraph")


def test_extract_head_from_pdf(monkeypatch, tmp_path):
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n")

    class DummyPage:
        def __init__(self, text: str):
            self._text = text

        def extract_text(self):
            return self._text

    class DummyReader:
        def __init__(self, *_):
            self.pages = [DummyPage("PDF snippet one"), DummyPage("Second page text")]

    monkeypatch.setattr("pypdf.PdfReader", lambda path: DummyReader(path))

    snippet = extract_head(pdf_path)

    assert snippet.startswith("PDF snippet one Second page text")


def test_extract_head_from_pptx(tmp_path):
    ppt_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Body paragraph"
    presentation.save(ppt_path)

    snippet = extract_head(ppt_path)

    assert snippet.startswith("Slide Title Body paragraph")


def test_extract_full_chunks_from_docx(tmp_path):
    doc_path = tmp_path / "long.docx"
    document = Document()
    for idx in range(10):
        document.add_paragraph(f"Paragraph {idx} " + "text " * 5)
    document.save(doc_path)

    chunks = extract_full_chunks(doc_path, chunk_size=50, overlap=0)

    assert chunks
    assert any("Paragraph" in chunk for chunk in chunks)


def test_extract_full_chunks_from_pdf(monkeypatch, tmp_path):
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n")

    class DummyPage:
        def __init__(self, text: str):
            self._text = text

        def extract_text(self):
            return self._text

    class DummyReader:
        def __init__(self, *_):
            self.pages = [DummyPage("One two three four five six seven eight nine ten" * 5)]

    monkeypatch.setattr("pypdf.PdfReader", lambda path: DummyReader(path))

    chunks = extract_full_chunks(pdf_path, chunk_size=40, overlap=0)

    assert len(chunks) >= 2


def test_extract_full_chunks_from_pptx(tmp_path):
    ppt_path = tmp_path / "chunks.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Chunk Title"
    slide.placeholders[1].text = "Paragraph one text"
    presentation.save(ppt_path)

    chunks = extract_full_chunks(ppt_path, chunk_size=20, overlap=0)

    assert chunks
    assert chunks[0].startswith("Chunk Title")


def test_extract_full_chunks_returns_empty_for_unknown_or_empty(tmp_path, monkeypatch):
    unknown = tmp_path / "data.bin"
    unknown.write_bytes(b"\x00\x01")
    assert extract_full_chunks(unknown) == []

    empty = tmp_path / "empty.txt"
    empty.write_text("", encoding="utf-8")
    assert extract_full_chunks(empty) == []

    # from_path failure should return no chunks for non-UTF8 payloads
    monkeypatch.setattr(
        "charset_normalizer.from_path",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(Exception("boom")),
    )
    text_path = tmp_path / "boom.txt"
    text_path.write_bytes(b"\xff\xfe\xfd")
    assert extract_full_chunks(text_path) == []


def test_extract_code_chunks_from_python(tmp_path):
    py_path = tmp_path / "sample.py"
    py_path.write_text(
        """\"\"\"Module docstring.\"\"\"

import os
import sys

CONSTANT = 1

def foo(a, b):
    \"\"\"Foo does bar.\"\"\"
    return a + b


class Bar:
    \"\"\"Bar class.\"\"\"

    def method(self, x):
        return x * 2

    async def async_method(self):
        return 42


TAIL_CONSTANT = "tail"
if __name__ == "__main__":
    print(TAIL_CONSTANT)
"""
    )

    chunks = extract_code_chunks(py_path)

    assert chunks
    assert [chunk.kind for chunk in chunks[:3]] == ["module", "function", "class"]
    assert chunks[1].name == "foo"
    assert "def foo" in chunks[1].display
    assert "Foo does bar" in chunks[1].text
    assert chunks[2].name == "Bar"
    assert "Methods: method, async_method" in chunks[2].text
    assert any(chunk.name == "Bar.method" for chunk in chunks)
    assert any(chunk.name == "Bar.async_method" for chunk in chunks)
    assert any("TAIL_CONSTANT" in chunk.text for chunk in chunks if chunk.kind == "module")


def test_extract_code_chunks_syntax_error_returns_empty(tmp_path):
    py_path = tmp_path / "broken.py"
    py_path.write_text("def nope(:\n    pass\n")

    assert extract_code_chunks(py_path) == []


def test_extract_code_chunks_non_python_returns_empty(tmp_path):
    text_path = tmp_path / "sample.txt"
    text_path.write_text("Hello\n")

    assert extract_code_chunks(text_path) == []


def test_extract_code_chunks_includes_leading_comments(tmp_path):
    py_path = tmp_path / "sample.py"
    py_path.write_text(
        """# Adds two numbers
# Returns the sum
def add(a, b):
    return a + b
"""
    )

    chunks = extract_code_chunks(py_path)
    func_chunk = next(chunk for chunk in chunks if chunk.kind == "function")
    assert "# Adds two numbers" in func_chunk.text
    assert func_chunk.start_line == 1


def test_extract_outline_chunks_from_markdown(tmp_path):
    md_path = tmp_path / "doc.md"
    md_path.write_text(
        """---
title: Demo
---

Intro before headings.

# Top
Top body.

## Child
Child body.

```python
### Not a heading
```

## Another
Another body.
"""
    )

    chunks = extract_outline_chunks(md_path, context_char_limit=200)

    assert chunks
    assert chunks[0].breadcrumb == "preamble"
    assert chunks[1].breadcrumb == "Top"
    assert chunks[2].breadcrumb == "Top > Child"
    assert chunks[3].breadcrumb == "Top > Another"
    assert not any(chunk.title == "Not a heading" for chunk in chunks)


def test_extract_full_chunks_with_lines_from_text(tmp_path):
    text_path = tmp_path / "sample.txt"
    text_path.write_text("a\nb\nc\nd\ne\n")

    chunks = extract_full_chunks_with_lines(text_path, chunk_size=100, overlap=0)

    assert chunks
    assert chunks[0].start_line == 1
    assert chunks[0].end_line == 5


def test_extract_full_chunks_with_lines_multiple_windows(tmp_path):
    text_path = tmp_path / "many.txt"
    text_path.write_text("L1\nL2\nL3\nL4\nL5\nL6\n", encoding="utf-8")

    chunks = extract_full_chunks_with_lines(text_path, chunk_size=4, overlap=0)

    assert len(chunks) >= 2
    assert all(chunk.start_line is not None and chunk.end_line is not None for chunk in chunks)
    assert chunks[0].start_line == 1
    assert chunks[0].end_line >= 1
    assert chunks[-1].end_line == 6


def test_extract_full_chunks_with_lines_non_text_has_no_lines(tmp_path):
    doc_path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Doc paragraph one")
    document.add_paragraph("Doc paragraph two")
    document.save(doc_path)

    chunks = extract_full_chunks_with_lines(doc_path, chunk_size=50, overlap=0)
    assert chunks
    assert chunks[0].start_line is None
    assert chunks[0].end_line is None


def test_extract_outline_chunks_no_headings_returns_empty(tmp_path):
    md_path = tmp_path / "plain.md"
    md_path.write_text("Just text.\nNo headings here.\n")

    assert extract_outline_chunks(md_path) == []


def test_extract_outline_chunks_non_markdown_returns_empty(tmp_path):
    txt_path = tmp_path / "doc.txt"
    txt_path.write_text("# Not markdown by extension\n")

    assert extract_outline_chunks(txt_path) == []


def test_extract_outline_chunks_supports_setext_headings(tmp_path):
    md_path = tmp_path / "setext.md"
    md_path.write_text(
        """Title
-----

Body line one.

Subtitle
=======

Body line two.
""",
        encoding="utf-8",
    )
    chunks = extract_outline_chunks(md_path, context_char_limit=200)
    assert chunks
    assert any(chunk.breadcrumb == "Title" for chunk in chunks)
    assert any(chunk.breadcrumb == "Subtitle" for chunk in chunks)


def test_extract_outline_chunks_ignores_headings_inside_tilde_fence(tmp_path):
    md_path = tmp_path / "fences.md"
    md_path.write_text(
        """# Outside

~~~python
## Not a heading
~~~

## Inside
Text.
""",
        encoding="utf-8",
    )
    chunks = extract_outline_chunks(md_path, context_char_limit=200)
    assert chunks
    assert any(chunk.breadcrumb == "Outside" for chunk in chunks)
    assert any(chunk.breadcrumb == "Outside > Inside" for chunk in chunks)
    assert not any(chunk.title == "Not a heading" for chunk in chunks)


def test_extract_code_chunks_accounts_for_decorators_and_no_symbols(tmp_path):
    decorated = tmp_path / "decorated.py"
    decorated.write_text(
        """@decorator
def foo():
    return 1
""",
        encoding="utf-8",
    )
    chunks = extract_code_chunks(decorated)
    assert chunks
    assert chunks[0].kind == "function"
    assert chunks[0].start_line == 1

    no_symbols = tmp_path / "nosymbols.py"
    no_symbols.write_text("X = 1\nY = 2\n", encoding="utf-8")
    chunks2 = extract_code_chunks(no_symbols)
    assert chunks2
    assert chunks2[0].kind == "module"


def test_extract_code_chunks_without_end_lineno_falls_back(tmp_path, monkeypatch):
    py_path = tmp_path / "no_end.py"
    py_path.write_text(
        """def foo():\n    x = 1\n    return x\n""",
        encoding="utf-8",
    )

    real_parse = ces.ast.parse

    def fake_parse(source: str):
        module = real_parse(source)
        for node in getattr(module, "body", []) or []:
            if hasattr(node, "end_lineno"):
                try:
                    delattr(node, "end_lineno")
                except Exception:
                    pass
            body = getattr(node, "body", None) or []
            for child in body:
                if hasattr(child, "end_lineno"):
                    try:
                        delattr(child, "end_lineno")
                    except Exception:
                        pass
        return module

    monkeypatch.setattr(ces.ast, "parse", fake_parse)
    chunks = extract_code_chunks(py_path)
    assert chunks
    assert chunks[0].name == "foo"


def test_text_readers_fallback_to_charset_detection(monkeypatch, tmp_path):
    data_path = tmp_path / "latin.txt"
    data_path.write_bytes(b"\xff\xfe")

    assert ces._read_text_head(data_path) is None
    assert ces._read_text_full(data_path) is None

    class EmptyResult:
        def __len__(self):
            return 0

    monkeypatch.setattr("charset_normalizer.from_path", lambda _path: EmptyResult())
    assert ces._read_text_head(data_path) is None
    assert ces._read_text_full(data_path) is None

    class NoBest:
        def __len__(self):
            return 1

        def best(self):
            return None

    monkeypatch.setattr("charset_normalizer.from_path", lambda _path: NoBest())
    assert ces._read_text_head(data_path) is None
    assert ces._read_text_full(data_path) is None

    class EmptyBest:
        def __len__(self):
            return 1

        def best(self):
            return ""

    monkeypatch.setattr("charset_normalizer.from_path", lambda _path: EmptyBest())
    assert ces._read_text_head(data_path) is None
    assert ces._read_text_full(data_path) is None

    class GoodBest:
        def __len__(self):
            return 1

        def best(self):
            return " Alpha \n Beta "

    monkeypatch.setattr("charset_normalizer.from_path", lambda _path: GoodBest())
    assert ces._read_text_head(data_path, char_limit=20) == "Alpha Beta"
    assert ces._read_text_full(data_path, char_limit=6) == " Alpha"
    assert ces._read_text_full(data_path, char_limit=0) == " Alpha \n Beta "


def test_read_text_utf8_oserror_and_zero_limit(tmp_path):
    missing = tmp_path / "missing.txt"
    assert ces._read_text_utf8(missing, char_limit=10) is None

    text_path = tmp_path / "sample.txt"
    text_path.write_text("abcdef", encoding="utf-8")
    assert ces._read_text_utf8(text_path, char_limit=0) == "abcdef"
    assert ces._read_text_utf8(text_path, char_limit=3) == "abc"


def test_pdf_extractor_error_and_empty_paths(monkeypatch, tmp_path):
    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF")

    monkeypatch.setattr(
        "pypdf.PdfReader",
        lambda _path: (_ for _ in ()).throw(RuntimeError("bad pdf")),
    )
    assert ces._pdf_extractor(pdf_path) is None

    class ErrorPage:
        def extract_text(self):
            raise RuntimeError("bad page")

    class BlankPage:
        def extract_text(self):
            return "   "

    class Reader:
        pages = [ErrorPage(), BlankPage()]

    monkeypatch.setattr("pypdf.PdfReader", lambda _path: Reader())
    assert ces._pdf_extractor(pdf_path) is None


def test_docx_extractor_error_and_empty(monkeypatch, tmp_path):
    doc_path = tmp_path / "sample.docx"
    doc_path.write_bytes(b"bad")

    monkeypatch.setattr(
        "docx.Document",
        lambda _path: (_ for _ in ()).throw(RuntimeError("bad docx")),
    )
    assert ces._docx_extractor(doc_path) is None

    monkeypatch.setattr("docx.Document", lambda _path: SimpleNamespace(paragraphs=[]))
    assert ces._docx_extractor(doc_path) is None

    monkeypatch.setattr(
        "docx.Document",
        lambda _path: SimpleNamespace(
            paragraphs=[
                SimpleNamespace(text="  "),
                SimpleNamespace(text="Alpha"),
                SimpleNamespace(text="Beta"),
            ]
        ),
    )
    assert ces._docx_extractor(doc_path, char_limit=6) == "Alpha "


def test_pptx_extractor_error_empty_and_shape_text(monkeypatch, tmp_path):
    ppt_path = tmp_path / "sample.pptx"
    ppt_path.write_bytes(b"bad")

    monkeypatch.setattr(
        "pptx.Presentation",
        lambda _path: (_ for _ in ()).throw(RuntimeError("bad pptx")),
    )
    assert ces._pptx_extractor(ppt_path) is None

    monkeypatch.setattr("pptx.Presentation", lambda _path: SimpleNamespace(slides=[]))
    assert ces._pptx_extractor(ppt_path) is None

    shape_with_text = SimpleNamespace(text_frame=None, text=" Fallback text ")
    paragraph_with_runs = SimpleNamespace(
        runs=[SimpleNamespace(text="Run "), SimpleNamespace(text="Text")],
        text="ignored",
    )
    paragraph_without_runs = SimpleNamespace(runs=[], text=" Paragraph text ")
    shape_with_frame = SimpleNamespace(
        text_frame=SimpleNamespace(
            paragraphs=[
                SimpleNamespace(runs=[], text=" "),
                paragraph_with_runs,
                paragraph_without_runs,
            ]
        )
    )
    empty_shape = SimpleNamespace(
        text_frame=SimpleNamespace(paragraphs=[SimpleNamespace(runs=[], text=" ")])
    )
    slide = SimpleNamespace(shapes=[empty_shape, shape_with_text, shape_with_frame])
    monkeypatch.setattr(
        "pptx.Presentation",
        lambda _path: SimpleNamespace(slides=[slide]),
    )

    extracted = ces._pptx_extractor(ppt_path, char_limit=100)

    assert extracted == "Fallback text Run Text Paragraph text"
    assert ces._extract_shape_text(empty_shape) is None
    assert ces._extract_shape_text(SimpleNamespace(text_frame=None, text=" ")) is None


def test_extract_full_chunks_with_lines_empty_and_non_text(monkeypatch, tmp_path):
    unknown = tmp_path / "data.bin"
    unknown.write_bytes(b"data")
    assert extract_full_chunks_with_lines(unknown) == []

    empty = tmp_path / "empty.txt"
    empty.write_text("", encoding="utf-8")
    assert extract_full_chunks_with_lines(empty) == []

    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF")
    monkeypatch.setattr(ces, "_pdf_extractor", lambda *_args, **_kwargs: "PDF text")
    chunks = extract_full_chunks_with_lines(pdf_path, chunk_size=10, overlap=0)
    assert chunks
    assert chunks[0].start_line is None


def test_outline_parser_edge_cases(tmp_path):
    md_path = tmp_path / "edge.md"
    md_path.write_text(
        (
            """---
title: Demo
---
# ignored in front matter window

"""
            + "###   \n"
            + """

# Title ###
Body

# Already heading
---

````
## ignored
```
still fenced
````

## Child
"""
            + ("x" * 50)
        ),
        encoding="utf-8",
    )

    chunks = extract_outline_chunks(md_path, context_char_limit=10)

    assert chunks
    assert any(chunk.title == "Title" for chunk in chunks)
    assert not any(chunk.title == "ignored" for chunk in chunks)
    assert all(len(chunk.text) <= 10 for chunk in chunks)
