"""Helpers to extract head snippets from various file types."""

from __future__ import annotations

import ast
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Protocol

from charset_normalizer import from_path
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

HEAD_CHAR_LIMIT = 1000
FULL_CHAR_LIMIT = 200_000
DEFAULT_CHUNK_SIZE = 1000
DEFAULT_CHUNK_OVERLAP = 100


class HeadExtractor(Protocol):
    """Protocol describing a file head extractor."""

    def __call__(self, path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
        ...


@dataclass(frozen=True)
class ExtractorEntry:
    extensions: tuple[str, ...]
    extractor: HeadExtractor


@dataclass(frozen=True, slots=True)
class CodeChunk:
    kind: str
    name: str
    display: str
    text: str
    start_line: int
    end_line: int


@dataclass(frozen=True, slots=True)
class OutlineChunk:
    level: int
    title: str
    breadcrumb: str
    text: str
    start_line: int
    end_line: int


@dataclass(frozen=True, slots=True)
class FullChunk:
    text: str
    start_line: int | None
    end_line: int | None


_registry: Dict[str, HeadExtractor] = {}

TEXT_EXTENSIONS = (
    ".txt",
    ".md",
    ".py",
    ".js",
    ".ts",
    ".json",
    ".yaml",
    ".yml",
    ".html",
    ".htm",
    ".toml",
    ".csv",
    ".log",
    ".ini",
    ".cfg",
    ".rst",
    ".tex",
    ".xml",
    ".sh",
    ".bat",
    ".go",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".rb",
    ".php",
    ".swift",
    ".rs",
    ".kt",
    ".dart",
    ".scala",
    ".pl",
    ".r",
    ".jl",
    ".hs",
    ".lua",
    ".vb",
    ".ps1",
    ".bash",
)


def register_extractor(entry: ExtractorEntry) -> None:
    for ext in entry.extensions:
        _registry[ext.lower()] = entry.extractor


def extract_head(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    """Return a text snippet representing the head of *path*."""

    extractor = _registry.get(path.suffix.lower())
    if extractor is None:
        return None
    return extractor(path, char_limit)


def extract_full_chunks(
    path: Path,
    *,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
    char_limit: int = FULL_CHAR_LIMIT,
) -> list[str]:
    """Return sliding-window chunks for text-like files."""

    suffix = path.suffix.lower()
    text: str | None = None
    if suffix in TEXT_EXTENSIONS:
        text = _read_text_full(path, char_limit)
    elif suffix == ".pdf":
        text = _pdf_extractor(path, char_limit)
    elif suffix == ".docx":
        text = _docx_extractor(path, char_limit)
    elif suffix == ".pptx":
        text = _pptx_extractor(path, char_limit)
    else:
        return []
    if text is None:
        return []
    normalized = text.replace("\r\n", "\n").strip()
    if not normalized:
        return []
    size = max(int(chunk_size), 1)
    stride = max(size - max(int(overlap), 0), 1)
    chunks: list[str] = []
    start = 0
    length = len(normalized)
    while start < length:
        window = normalized[start : start + size].strip()
        if window:
            chunks.append(window)
        if start + size >= length:
            break
        start += stride
    return chunks


def extract_full_chunks_with_lines(
    path: Path,
    *,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
    char_limit: int = FULL_CHAR_LIMIT,
) -> list[FullChunk]:
    """Return sliding-window chunks and approximate line ranges for text-like files.

    Line ranges are computed only for plain text inputs (TEXT_EXTENSIONS). Other extractors
    return chunks with line metadata set to None.
    """

    suffix = path.suffix.lower()
    text: str | None = None
    include_lines = False
    if suffix in TEXT_EXTENSIONS:
        include_lines = True
        text = _read_text_full(path, char_limit)
    elif suffix == ".pdf":
        text = _pdf_extractor(path, char_limit)
    elif suffix == ".docx":
        text = _docx_extractor(path, char_limit)
    elif suffix == ".pptx":
        text = _pptx_extractor(path, char_limit)
    else:
        return []
    if text is None:
        return []

    normalized = text.replace("\r\n", "\n").strip()
    if not normalized:
        return []

    size = max(int(chunk_size), 1)
    stride = max(size - max(int(overlap), 0), 1)
    chunks: list[FullChunk] = []
    start = 0
    length = len(normalized)
    while start < length:
        end = min(start + size, length)
        window = normalized[start:end]
        cleaned = window.strip()
        if cleaned:
            start_line: int | None = None
            end_line: int | None = None
            if include_lines:
                leading = len(window) - len(window.lstrip())
                trailing = len(window) - len(window.rstrip())
                span_start = min(start + leading, length)
                span_end = max(span_start, end - trailing)
                start_line = normalized.count("\n", 0, span_start) + 1
                last_index = max(span_start, span_end - 1)
                end_line = normalized.count("\n", 0, last_index) + 1
            chunks.append(FullChunk(text=cleaned, start_line=start_line, end_line=end_line))
        if end >= length:
            break
        start += stride
    return chunks


def extract_code_chunks(
    path: Path,
    *,
    char_limit: int = FULL_CHAR_LIMIT,
) -> list[CodeChunk]:
    """Return AST-aware code chunks for supported languages (Python only for now)."""

    if path.suffix.lower() != ".py":
        return []

    source = _read_text_full(path, char_limit)
    if not source:
        return []
    source = source.replace("\r\n", "\n")

    try:
        module = ast.parse(source)
    except SyntaxError:
        return []

    lines = source.splitlines(keepends=True)
    max_line = len(lines)

    def clamp_line(value: int) -> int:
        if value < 1:
            return 1
        if value > max_line:
            return max_line
        return value

    def node_start_line(node) -> int:
        lineno = getattr(node, "lineno", None)
        start = int(lineno) if isinstance(lineno, int) else 1
        decorators = getattr(node, "decorator_list", None) or []
        for deco in decorators:
            deco_line = getattr(deco, "lineno", None)
            if isinstance(deco_line, int):
                start = min(start, deco_line)
        return clamp_line(start)

    def node_end_line(node) -> int:
        end_lineno = getattr(node, "end_lineno", None)
        if isinstance(end_lineno, int):
            return clamp_line(end_lineno)
        body = getattr(node, "body", None) or []
        if body:
            last = body[-1]
            last_end = getattr(last, "end_lineno", None)
            if isinstance(last_end, int):
                return clamp_line(last_end)
            last_line = getattr(last, "lineno", None)
            if isinstance(last_line, int):
                return clamp_line(last_line)
        lineno = getattr(node, "lineno", None)
        if isinstance(lineno, int):
            return clamp_line(lineno)
        return max_line

    def slice_lines(start: int, end: int) -> str:
        if not max_line:
            return ""
        start = clamp_line(start)
        end = clamp_line(end)
        if end < start:
            end = start
        return "".join(lines[start - 1 : end]).strip()

    def signature_line(node) -> str:
        lineno = getattr(node, "lineno", None)
        if not isinstance(lineno, int):
            return ""
        idx = lineno - 1
        if idx < 0 or idx >= max_line:
            return ""
        return lines[idx].strip()

    chunks: list[CodeChunk] = []

    symbols: list[tuple[int, int, object]] = []
    for node in module.body:
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
            symbols.append((node_start_line(node), node_end_line(node), node))
    symbols.sort(key=lambda item: item[0])

    def add_module_chunk(start: int, end: int, *, prelude: bool) -> None:
        text = slice_lines(start, end)
        if not text:
            return
        chunks.append(
            CodeChunk(
                kind="module",
                name="module" if prelude else "module_globals",
                display="module" if prelude else "module globals",
                text=text,
                start_line=start,
                end_line=end,
            )
        )

    if not symbols:
        add_module_chunk(1, max_line, prelude=True)
        return chunks

    cursor = 1
    seen_symbol = False

    for start, end, node in symbols:
        if cursor <= start - 1:
            add_module_chunk(cursor, start - 1, prelude=not seen_symbol)
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            text = slice_lines(start, end)
            if text:
                chunks.append(
                    CodeChunk(
                        kind="function",
                        name=node.name,
                        display=signature_line(node) or f"def {node.name}",
                        text=text,
                        start_line=start,
                        end_line=end,
                    )
                )
            cursor = end + 1
            seen_symbol = True
            continue

        if isinstance(node, ast.ClassDef):
            class_display = signature_line(node) or f"class {node.name}"
            docstring = ast.get_docstring(node) or ""
            method_names: list[str] = []
            for child in node.body:
                if isinstance(child, (ast.FunctionDef, ast.AsyncFunctionDef)):
                    method_names.append(child.name)

            def is_docstring_expr(stmt: object) -> bool:
                if not isinstance(stmt, ast.Expr):
                    return False
                value = getattr(stmt, "value", None)
                if isinstance(value, ast.Constant):
                    return isinstance(value.value, str)
                return False

            class_parts = [slice_lines(start, node.lineno)]
            if docstring.strip():
                class_parts.append(docstring.strip())
            for idx, child in enumerate(node.body):
                if isinstance(child, (ast.FunctionDef, ast.AsyncFunctionDef)):
                    continue
                if idx == 0 and docstring and is_docstring_expr(child):
                    continue
                child_text = slice_lines(node_start_line(child), node_end_line(child))
                if child_text:
                    class_parts.append(child_text)
            if method_names:
                class_parts.append("Methods: " + ", ".join(method_names))
            class_text = "\n".join(part for part in class_parts if part).strip()
            if class_text:
                chunks.append(
                    CodeChunk(
                        kind="class",
                        name=node.name,
                        display=class_display,
                        text=class_text,
                        start_line=start,
                        end_line=end,
                    )
                )

            for child in node.body:
                if not isinstance(child, (ast.FunctionDef, ast.AsyncFunctionDef)):
                    continue
                child_start = node_start_line(child)
                child_end = node_end_line(child)
                text = slice_lines(child_start, child_end)
                if not text:
                    continue
                raw_sig = signature_line(child)
                display = f"{node.name}.{child.name}"
                if raw_sig:
                    normalized_sig = raw_sig.strip()
                    if normalized_sig.startswith("async def "):
                        tail = normalized_sig[len("async def ") :].rstrip(":").strip()
                        display = f"async {node.name}.{tail}"
                    elif normalized_sig.startswith("def "):
                        tail = normalized_sig[len("def ") :].rstrip(":").strip()
                        display = f"{node.name}.{tail}"
                    else:
                        display = f"{node.name}.{normalized_sig.rstrip(':').strip()}"
                chunks.append(
                    CodeChunk(
                        kind="method",
                        name=f"{node.name}.{child.name}",
                        display=display,
                        text=text,
                        start_line=child_start,
                        end_line=child_end,
                    )
                )

            cursor = end + 1
            seen_symbol = True
            continue

        cursor = end + 1
        seen_symbol = True

    if cursor <= max_line:
        add_module_chunk(cursor, max_line, prelude=False)

    return chunks


def extract_outline_chunks(
    path: Path,
    *,
    context_char_limit: int = 800,
    char_limit: int = FULL_CHAR_LIMIT,
) -> list[OutlineChunk]:
    """Return outline chunks for Markdown files (headings + section snippets)."""

    if path.suffix.lower() not in {".md", ".markdown", ".mdx"}:
        return []

    source = _read_text_full(path, char_limit)
    if not source:
        return []
    source = source.replace("\r\n", "\n")
    lines = source.splitlines()
    if not lines:
        return []

    front_matter_end: int | None = None
    if lines and lines[0].strip() == "---":
        for idx, line in enumerate(lines[1:], start=1):
            marker = line.strip()
            if marker in {"---", "..."}:
                front_matter_end = idx
                break

    @dataclass(frozen=True, slots=True)
    class Heading:
        line: int
        end_line: int
        level: int
        title: str
        content_start: int

    headings: list[Heading] = []
    heading_starts: set[int] = set()

    in_fence = False
    fence_char = ""
    fence_len = 0
    prev_line_text: str | None = None
    prev_line_index: int | None = None

    fence_re = re.compile(r"^\s*([`~]{3,})")
    atx_re = re.compile(r"^\s{0,3}(#{1,6})\s+(.*)$")
    setext_re = re.compile(r"^\s{0,3}([=-]{3,})\s*$")

    def handle_fence(line: str) -> bool:
        nonlocal in_fence, fence_char, fence_len, prev_line_text, prev_line_index
        match = fence_re.match(line)
        if not match:
            return False
        marker = match.group(1)
        if not in_fence:
            in_fence = True
            fence_char = marker[0]
            fence_len = len(marker)
            prev_line_text = None
            prev_line_index = None
            return True
        if marker[0] == fence_char and len(marker) >= fence_len:
            in_fence = False
            fence_char = ""
            fence_len = 0
            prev_line_text = None
            prev_line_index = None
            return True
        return True

    def record_atx(line_no: int, marker: str, raw_title: str) -> None:
        title = raw_title.strip()
        title = re.sub(r"\s#+\s*$", "", title).strip()
        if not title:
            return
        if line_no in heading_starts:
            return
        heading_starts.add(line_no)
        headings.append(
            Heading(
                line=line_no,
                end_line=line_no,
                level=len(marker),
                title=title,
                content_start=line_no + 1,
            )
        )

    def record_setext(title_line: int, underline_line: int, underline: str, title: str) -> None:
        title = title.strip()
        if not title:
            return
        if title_line in heading_starts:
            return
        heading_starts.add(title_line)
        level = 1 if underline.startswith("=") else 2
        headings.append(
            Heading(
                line=title_line,
                end_line=underline_line,
                level=level,
                title=title,
                content_start=underline_line + 1,
            )
        )

    for idx, line in enumerate(lines, start=1):
        if front_matter_end is not None and idx <= front_matter_end + 1:
            prev_line_text = None
            prev_line_index = None
            continue
        if handle_fence(line):
            continue
        if in_fence:
            continue

        atx = atx_re.match(line)
        if atx:
            record_atx(idx, atx.group(1), atx.group(2))
            prev_line_text = None
            prev_line_index = None
            continue

        setext = setext_re.match(line)
        if setext and prev_line_text and prev_line_index:
            underline = setext.group(1)
            if prev_line_text.lstrip().startswith("#"):
                prev_line_text = line
                prev_line_index = idx
                continue
            record_setext(prev_line_index, idx, underline, prev_line_text)
            prev_line_text = None
            prev_line_index = None
            continue

        prev_line_text = line
        prev_line_index = idx

    if not headings:
        return []

    headings.sort(key=lambda item: item.line)

    preamble_start = 1
    if front_matter_end is not None:
        preamble_start = front_matter_end + 2
    first_heading_line = headings[0].line
    if preamble_start <= first_heading_line - 1:
        preamble_text = _cleanup_snippet("\n".join(lines[preamble_start - 1 : first_heading_line - 1]))
        if preamble_text:
            headings.insert(
                0,
                Heading(
                    line=preamble_start,
                    end_line=first_heading_line - 1,
                    level=0,
                    title="preamble",
                    content_start=preamble_start,
                ),
            )

    stack: list[tuple[int, str]] = []
    chunks: list[OutlineChunk] = []

    for idx, heading in enumerate(headings):
        if heading.level == 0:
            breadcrumb = "preamble"
        else:
            while stack and stack[-1][0] >= heading.level:
                stack.pop()
            stack.append((heading.level, heading.title))
            breadcrumb = " > ".join(title for _, title in stack)

        section_end = len(lines)
        for next_heading in headings[idx + 1 :]:
            if next_heading.line <= heading.line:
                continue
            if heading.level == 0:
                section_end = next_heading.line - 1
                break
            if next_heading.level <= heading.level:
                section_end = next_heading.line - 1
                break

        start = max(heading.content_start, 1)
        end = max(section_end, start)
        section_text = "\n".join(lines[start - 1 : end])
        cleaned = _cleanup_snippet(section_text) or ""
        if context_char_limit > 0 and len(cleaned) > context_char_limit:
            cleaned = cleaned[:context_char_limit].rstrip()

        chunks.append(
            OutlineChunk(
                level=heading.level,
                title=heading.title,
                breadcrumb=breadcrumb,
                text=cleaned,
                start_line=heading.line,
                end_line=section_end,
            )
        )

    return chunks


# Placeholder extractors ----------------------------------------------------

def _read_text_head(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    """Return the first *char_limit* characters of a text-like file."""

    try:
        result = from_path(path)
    except Exception:
        return None
    if result is None or not len(result):
        return None
    best = result.best()
    if best is None:
        return None
    text = str(best)
    if not text:
        return None
    snippet = text[:char_limit]
    return _cleanup_snippet(snippet)


def _read_text_full(path: Path, char_limit: int = FULL_CHAR_LIMIT) -> str | None:
    """Return up to *char_limit* characters from a text-like file."""

    try:
        result = from_path(path)
    except Exception:
        return None
    if result is None or not len(result):
        return None
    best = result.best()
    if best is None:
        return None
    text = str(best)
    if not text:
        return None
    if char_limit > 0:
        return text[:char_limit]
    return text


def _pdf_extractor(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    try:
        reader = PdfReader(str(path))
    except Exception:
        return None
    buffer: list[str] = []
    total_chars = 0
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if not text:
            continue
        buffer.append(text)
        total_chars += len(text)
        if total_chars >= char_limit:
            break
    combined = "\n".join(buffer)
    if not combined:
        return None
    cleaned = _cleanup_snippet(combined)
    if not cleaned:
        return None
    return cleaned[:char_limit]


def _docx_extractor(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    try:
        document = Document(str(path))
    except Exception:
        return None
    buffer: list[str] = []
    total_chars = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        buffer.append(text)
        total_chars += len(text)
        if total_chars >= char_limit:
            break
    combined = "\n".join(buffer)
    if not combined:
        return None
    cleaned = _cleanup_snippet(combined)
    if not cleaned:
        return None
    return cleaned[:char_limit]


def _pptx_extractor(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    try:
        presentation = Presentation(str(path))
    except Exception:
        return None
    buffer: list[str] = []
    total_chars = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = _extract_shape_text(shape)
            if not text:
                continue
            buffer.append(text)
            total_chars += len(text)
            if total_chars >= char_limit:
                break
        if total_chars >= char_limit:
            break
    combined = "\n".join(buffer)
    if not combined:
        return None
    cleaned = _cleanup_snippet(combined)
    if not cleaned:
        return None
    return cleaned[:char_limit]


def _extract_shape_text(shape) -> str | None:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        text = getattr(shape, "text", "")
        text = text.strip()
        return text or None
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        if getattr(paragraph, "runs", None):
            text = "".join(run.text for run in paragraph.runs)
        else:
            text = paragraph.text
        text = (text or "").strip()
        if text:
            paragraphs.append(text)
    if not paragraphs:
        return None
    return " ".join(paragraphs)


def _cleanup_snippet(snippet: str) -> str | None:
    lines = [line.strip() for line in snippet.splitlines() if line.strip()]
    joined = " ".join(lines)
    return joined or None


def _unimplemented_extractor(path: Path, char_limit: int = HEAD_CHAR_LIMIT) -> str | None:
    return None


register_extractor(
    ExtractorEntry(
        extensions=TEXT_EXTENSIONS,
        extractor=_read_text_head,
    )
)

register_extractor(
    ExtractorEntry((".pdf",), _pdf_extractor)
)

register_extractor(
    ExtractorEntry((".docx",), _docx_extractor)
)

register_extractor(
    ExtractorEntry((".pptx",), _pptx_extractor)
)
